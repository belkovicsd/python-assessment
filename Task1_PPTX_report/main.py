from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import PIL.Image as Image
import logging
import matplotlib.pyplot
import numpy
import os
import pandas
import urllib3

CONTENT_KEY = 'content'
FILE_NAME = 'example_output.pptx'
PICTURE_URL = 'https://i0.wp.com/mechguru.com/wp-content/uploads/2012/05/S_N-Diagram.jpeg?w=688&ssl=1'

TABLE_OF_CONTENT_LAYOUT = 1
TITLE_LAYOUT = 0
TITLE_ONLY_LAYOUT = 5


class UnknownTypeException(Exception):
    pass


def generate_report():
    clear_presentation()

    conf_path = get_file_path('sample.json')
    file = pandas.read_json(conf_path)

    for index, row in file.iterrows():

        first_element = row.iloc[0]
        type_value = first_element['type']

        process_object = get_type_mapping(type_value)

        try:
            if process_object is None:
                raise UnknownTypeException("Invalid 'type' value:")

            else:
                logging.debug(f"read_config - Processing record, Type: '{type_value}'")
                process_object(first_element)

        except UnknownTypeException as e:
            logging.warning(f"read_config - {e} '{type_value}', Index: {index}")


def get_type_mapping(type_value):
    type_mapping = {
        'title': generate_title_slide_report,
        'text': generate_text_slide_report,
        'list': generate_list_slide_report,
        'picture': generate_picture_slide_report,
        'plot': generate_plot_slide_report
    }
    return type_mapping.get(type_value)


def get_file_path(filename):
    dir_path = os.path.dirname(os.path.realpath(__file__))
    return os.path.join(dir_path, filename)


def clear_presentation():
    Presentation().save(FILE_NAME)
    logging.debug('clear_presentation - Empty Presentation created')


def create_paragraph_for_each_element(elements, slide):
    shape = slide.shapes.placeholders[1]
    text_frame = shape.text_frame
    for element in elements:
        paragraph = text_frame.add_paragraph()
        paragraph.text = element['text']
        paragraph.level = element['level']


def create_presentation(layout, data):
    presentation = Presentation(FILE_NAME)
    slide_layout = presentation.slide_layouts[layout]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = pandas.Series(data)['title']

    return presentation, slide


def generate_title_slide_report(data):
    presentation, slide = create_presentation(TITLE_LAYOUT, data)

    subtitle = slide.placeholders[1]
    subtitle.text = pandas.Series(data)[CONTENT_KEY]

    presentation.save(FILE_NAME)


def generate_text_slide_report(data):
    presentation, slide = create_presentation(TITLE_ONLY_LAYOUT, data)

    text_box = slide.shapes.add_textbox(Inches(1.28), Inches(1.18), Inches(5.81), Inches(0.4))
    text_box.text_frame.text = pandas.Series(data)[CONTENT_KEY]

    presentation.save(FILE_NAME)


def generate_list_slide_report(data):
    presentation, slide = create_presentation(TABLE_OF_CONTENT_LAYOUT, data)

    elements = pandas.Series(data)[CONTENT_KEY]
    create_paragraph_for_each_element(elements, slide)

    presentation.save(FILE_NAME)


def generate_picture_slide_report(data):
    presentation, slide = create_presentation(TITLE_ONLY_LAYOUT, data)

    http = urllib3.PoolManager()
    response = http.request('GET', PICTURE_URL)

    if response.status != 200:
        logging.error(f"generate_picture_slide_report - Response status: '{response.status}'")

    else:
        try:
            downloaded_image = Image.open(BytesIO(response.data))
            downloaded_image.save('picture.png')
            slide.shapes.add_picture(
                pandas.Series(data)[CONTENT_KEY], Inches(1.42), Inches(1.21), Inches(7.185), Inches(4.54)
            )

        except Exception as e:
            logging.error(f"generate_picture_slide_report - An error occurred while processing the image: '{e}'")

        finally:
            if 'downloaded_image' in locals():
                downloaded_image.close()
                os.remove('picture.png')
                logging.debug('generate_picture_slide_report - The downloaded picture has been deleted!')

    presentation.save(FILE_NAME)


def generate_plot_slide_report(data):
    presentation, slide = create_presentation(TITLE_ONLY_LAYOUT, data)

    plot_data_path = get_file_path(pandas.Series(data)[CONTENT_KEY])
    values = numpy.loadtxt(plot_data_path, delimiter=';')
    x = [row[0] for row in values]
    y = [row[1] for row in values]
    configuration = pandas.Series(data)['configuration']
    matplotlib.pyplot.plot(x, y)
    matplotlib.pyplot.xlabel(configuration['x-label'])
    matplotlib.pyplot.ylabel(configuration['y-label'])

    buffer = BytesIO()
    matplotlib.pyplot.savefig(buffer, format='png')
    buffer.seek(0)
    matplotlib.pyplot.close()

    slide.shapes.add_picture(buffer, Inches(1.38), Inches(1.18), Inches(6.4), Inches(4.8))

    presentation.save(FILE_NAME)


if __name__ == "__main__":
    generate_report()
